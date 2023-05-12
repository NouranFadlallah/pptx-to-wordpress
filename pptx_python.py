import collections.abc
from pptx import Presentation

sourceFile = './HKSIP1E.pptx'
prs = Presentation(sourceFile)
speaker_notes = []
title_texts = []
for num, slide in enumerate(prs.slides):
    print(f"Slide no: {num}")
    if slide.has_notes_slide:
        speaker_notes.append(slide.notes_slide.notes_text_frame.text)
    for shape in slide.shapes:
        if shape.is_placeholder:
            if shape.placeholder_format.idx == 0:  # TITLE
                title_texts.append(shape.text)
                print(shape.text)
